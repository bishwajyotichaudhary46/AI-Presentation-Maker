import openai
from pptx.utill import Pt
import os
from pptx.dml.color import RGBColor
from dotenv import load_dotenv
from pptx import Presentation

load_dotenv()

openai.api_key = os.getenv("OPENAI_AI_KEY")

TITLE_FONT_SIZE = Pt(32)
CONTENT_FONT_SIZE = Pt(16)

def create_slide_titles(topic, num_slides):
    prompt =  "Generate {num_slides} short slide titles from the topic '{topic}' ."
    completion = openai.ChatCompetion.create(
        model = 'gpt-3.5-turbo',
        messages = [{'role': "system", "content":prompt}],
        temperature = 0.0,
        top_p = 0.1,
        max_tokens = 200,
        request_timeout = 15,
    )

    return completion.choices[0].message.content

def create_slide_content(slide_titles):
    prompt =  f"Generate content for the slide: {slide_titles}. The content must be in medium worded paragraphs. Only return 2 paragraphs."
    completion = openai.ChatCompetion.create(
        model = 'gpt-3.5-turbo',
        messages = [{'role': "system", "content":prompt}],
        temperature = 0.0,
        top_p = 0.1,
        max_tokens = 300,
        request_timeout = 15,
    )

    return completion.choices[0].message.content

def create_presentation(topic, slide_titles, slide_content):
    powerpoint = Presentation()

    title_slide_layout = powerpoint.slide_layouts[0]
    content_slide_layout = powerpoint.slide_layouts[1]

    background_color = RGBColor(173, 216, 230)

    title_slide = powerpoint.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = topic
    
    title.text_frame.paragraphs[0].front.size = TITLE_FONT_SIZE
    title.text_frame.paragraphs[0].front.bold = True
    content = title_slide.placeholders[1]
    content.text = "Created By AI."
    
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color

    for slide_title, slide_content in zip(slide_titles, slide_content):

        slide = powerpoint.slide.add_slide(content_slide_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = background_color

        title = slide.shapes.title
        title.text = slide_title
        title.text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
        title.text_frame.paragraphs[0].font.bold = True

        content = slide.placeholders[1]
        content = slide_content
        for paragraph in content.text_frame.paragraphs:
            paragraph.font.size = CONTENT_FONT_SIZE

    powerpoint.save(f"powerpoints/{topic}.pptx")

def main():
    topic = "AI in AutoFlight System"
    num_slides = 1

    slide_titles = create_slide_titles(topic, num_slides)
    print("Generated Slide Titles. ")
    filtered_slide_titles = [item for item in slide_titles if item.strip() != '']

    slide_contents = [create_slide_content(title) for title in [slide_titles]]
    print("Genrated slide contents.")
    create_presentation(topic, filtered_slide_titles, slide_contents)

if __name__ == '__main__':
    main()
    




